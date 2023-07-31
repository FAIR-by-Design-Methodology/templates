from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
import re
import glob
import yaml

DOI_PATTERN = r'\b10\.\d{4,}(?:\.\d+)*\/[^\s]+\b'

def replace_text_last_slide(presentation_path: str, new_doi: str):
    presentation = Presentation(presentation_path)
    last_slide = presentation.slides[-1]

    for shape in last_slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                doi_match = re.findall(DOI_PATTERN, paragraph.text)
                if 'To cite this presentation' in paragraph.text or (len(doi_match)>0 and doi_match[0]!=new_doi):
                    print(f'Found replacement candidate in {presentation_path}, replacing...')
                    paragraph.text = ''
                    run_instructions = paragraph.add_run()
                    run_instructions.text = 'Cite instructions on:\n'
                    run_instructions.font.name = 'Quicksand'
                    run_instructions.font.size = Pt(12)
                    run_instructions.font.color.rgb = RGBColor(0, 84, 144)
                    run_doi = paragraph.add_run()
                    run_doi.text = f'https://doi.org/{new_doi}'
                    run_doi.font.name = 'Quicksand'
                    run_doi.font.size = Pt(12)
                    run_doi.font.color.rgb = RGBColor(229, 115, 0)
                    presentation.save(presentation_path)
                    return

def get_doi():
    with open('CITATION.cff', 'r') as yaml_file:
        yaml_content = yaml.safe_load(yaml_file)
        return yaml_content.get('doi')

current_doi = get_doi()

for file in glob.glob('**/*.pptx', recursive=True):
    if 'venv/' in file:
        continue
    print(f'Checking {file}...')
    replace_text_last_slide(file, current_doi)
